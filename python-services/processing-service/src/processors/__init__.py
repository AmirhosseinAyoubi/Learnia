import logging
import pdfplumber
from pptx import Presentation
import tiktoken

logger = logging.getLogger(__name__)

_ENCODING = "cl100k_base"
_CHUNK_TOKENS = 512


def extract_text(file_url: str, file_type: str) -> tuple[str, int]:
    """Extract full text and page count from a document file."""
    ft = file_type.upper()
    if ft == "PDF":
        return _extract_pdf(file_url)
    elif ft in ("PPTX", "PPT"):
        return _extract_pptx(file_url)
    elif ft == "TXT":
        return _extract_txt(file_url)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def chunk_text(text: str, max_tokens: int = _CHUNK_TOKENS) -> list[str]:
    """Split text into token-bounded chunks using tiktoken."""
    if not text.strip():
        return []
    enc = tiktoken.get_encoding(_ENCODING)
    tokens = enc.encode(text)
    chunks = []
    for i in range(0, len(tokens), max_tokens):
        chunk_tokens = tokens[i : i + max_tokens]
        chunks.append(enc.decode(chunk_tokens))
    return chunks


def _extract_pdf(path: str) -> tuple[str, int]:
    parts = []
    with pdfplumber.open(path) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                parts.append(t)
    return "\n\n".join(parts), page_count


def _extract_pptx(path: str) -> tuple[str, int]:
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n\n".join(parts), len(prs.slides)


def _extract_txt(path: str) -> tuple[str, int]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    page_count = max(1, len(text.split("\n")) // 50)
    return text, page_count
