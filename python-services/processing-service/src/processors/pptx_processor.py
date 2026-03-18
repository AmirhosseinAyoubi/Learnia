"""
processors/pptx_processor.py — Extracts plain text from PowerPoint files.

Iterates over every slide and pulls text from all text frames and table cells.
"""

import io
from pptx import Presentation


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    Extract all text from a PPTX file.

    Args:
        file_bytes: raw bytes of the PPTX file

    Returns:
        Concatenated plain text from all slides.
        Each slide's content is separated by a blank line.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    slide_texts = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            # Text frames (titles, text boxes, content placeholders)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in paragraph.runs).strip()
                    if line:
                        parts.append(line)
            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)

        if parts:
            slide_texts.append(f"[Slide {slide_number}]\n" + "\n".join(parts))

    return "\n\n".join(slide_texts)
